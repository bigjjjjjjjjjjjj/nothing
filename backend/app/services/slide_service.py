"""講義檔案處理服務"""
import os
import io
from typing import Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)


class SlideProcessingError(Exception):
    """講義處理錯誤"""
    pass


class SlideService:
    """講義檔案處理服務"""

    SUPPORTED_EXTENSIONS = {
        '.pdf': 'application/pdf',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.doc': 'application/msword',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    }

    def __init__(self, upload_dir: str = "uploads"):
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)

    async def process_file(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        處理上傳的檔案

        Args:
            file_content: 檔案內容 (bytes)
            filename: 檔案名稱

        Returns:
            包含文字內容、頁數等資訊的字典
        """
        file_ext = os.path.splitext(filename)[1].lower()

        if file_ext not in self.SUPPORTED_EXTENSIONS:
            raise SlideProcessingError(f"不支援的檔案格式: {file_ext}")

        try:
            if file_ext == '.pdf':
                return await self._process_pdf(file_content, filename)
            elif file_ext in ['.ppt', '.pptx']:
                return await self._process_powerpoint(file_content, filename)
            elif file_ext in ['.doc', '.docx']:
                return await self._process_word(file_content, filename)
        except Exception as e:
            logger.error(f"處理檔案失敗: {filename}, 錯誤: {str(e)}")
            raise SlideProcessingError(f"處理檔案失敗: {str(e)}")

    async def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """處理 PDF 檔案"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            total_pages = len(pdf_reader.pages)
            extracted_text = ""

            for page_num in range(total_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                extracted_text += f"\n--- 第 {page_num + 1} 頁 ---\n{text}\n"

            return {
                'filename': filename,
                'total_pages': total_pages,
                'extracted_text': extracted_text.strip(),
                'file_type': 'pdf'
            }

        except Exception as e:
            raise SlideProcessingError(f"PDF 處理失敗: {str(e)}")

    async def _process_powerpoint(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """處理 PowerPoint 檔案"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)

            total_slides = len(presentation.slides)
            extracted_text = ""

            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_text = f"\n--- 投影片 {slide_num} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                    # 處理表格
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            slide_text += row_text + "\n"

                extracted_text += slide_text

            return {
                'filename': filename,
                'total_pages': total_slides,
                'extracted_text': extracted_text.strip(),
                'file_type': 'powerpoint'
            }

        except Exception as e:
            raise SlideProcessingError(f"PowerPoint 處理失敗: {str(e)}")

    async def _process_word(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """處理 Word 檔案"""
        try:
            doc_file = io.BytesIO(file_content)
            document = Document(doc_file)

            extracted_text = ""
            total_paragraphs = len(document.paragraphs)

            for para_num, paragraph in enumerate(document.paragraphs, start=1):
                if paragraph.text.strip():
                    extracted_text += paragraph.text + "\n"

            # 處理表格
            for table in document.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    extracted_text += row_text + "\n"

            # Word 文件沒有明確的「頁數」概念，使用段落數估計
            estimated_pages = max(1, total_paragraphs // 10)

            return {
                'filename': filename,
                'total_pages': estimated_pages,
                'extracted_text': extracted_text.strip(),
                'file_type': 'word'
            }

        except Exception as e:
            raise SlideProcessingError(f"Word 處理失敗: {str(e)}")

    async def save_file(self, file_content: bytes, filename: str) -> str:
        """儲存檔案到本地"""
        file_path = os.path.join(self.upload_dir, filename)

        try:
            with open(file_path, 'wb') as f:
                f.write(file_content)
            return file_path
        except Exception as e:
            logger.error(f"儲存檔案失敗: {filename}, 錯誤: {str(e)}")
            raise SlideProcessingError(f"儲存檔案失敗: {str(e)}")

    def is_supported_file(self, filename: str) -> bool:
        """檢查檔案格式是否支援"""
        file_ext = os.path.splitext(filename)[1].lower()
        return file_ext in self.SUPPORTED_EXTENSIONS


# 建立全域實例
slide_service = SlideService()
